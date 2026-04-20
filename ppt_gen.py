"""PPT 생성 모듈 — BCG/Deloitte 컨설팅 스타일 (레이아웃 개선판)"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import io, datetime

NAVY  = RGBColor(0x1E,0x3A,0x5F)
CYAN  = RGBColor(0x00,0xB4,0xD8)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LGRAY = RGBColor(0xB0,0xC4,0xD8)
DGRAY = RGBColor(0x0F,0x1F,0x35)
CYAN2 = RGBColor(0x00,0x8B,0xAA)
W, H  = Inches(13.33), Inches(7.5)

def _prs():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    return prs

def _bg(slide, color=NAVY):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def _rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _txt(slide, x, y, w, h, text, size=13, bold=False,
         color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    run = p.add_run(); run.text = str(text)
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = color; run.font.name = "Calibri"
    return tb

def _header(slide, title, sub=""):
    _rect(slide, 0, 0, 13.33, 1.0, DGRAY)
    _rect(slide, 0, 0, 0.07, 1.0, CYAN)
    _txt(slide, 0.18, 0.1, 10, 0.5, title, 24, True, WHITE)
    if sub:
        _txt(slide, 0.18, 0.6, 10, 0.35, sub, 11, False, LGRAY)
    today = datetime.date.today().strftime("%Y.%m.%d")
    _txt(slide, 10.8, 0.35, 2.4, 0.3, today, 10, False, LGRAY, PP_ALIGN.RIGHT)

def _footer(slide):
    _rect(slide, 0, 7.15, 13.33, 0.35, DGRAY)
    _txt(slide, 0.3, 7.18, 8, 0.28,
         "AI NPD SUITE — 네이버 음료 트렌드 분석", 9, False, LGRAY)
    _txt(slide, 10.5, 7.18, 2.6, 0.28,
         "CONFIDENTIAL", 9, False, LGRAY, PP_ALIGN.RIGHT)

# ── 슬라이드 1: 타이틀
def slide_title(prs, keywords, n_blog, n_shop):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _rect(s, 0, 0, 0.45, 7.5, DGRAY)
    _rect(s, 0.45, 5.3, 12.88, 2.2, DGRAY)
    _rect(s, 0.45, 5.28, 12.88, 0.05, CYAN)
    _txt(s, 1.2, 1.4, 11, 0.55, "🍹  AI NPD SUITE", 18, True, CYAN)
    _txt(s, 1.2, 1.95, 11, 1.1, "네이버 음료 트렌드\n분석 리포트", 34, True, WHITE)
    _txt(s, 1.2, 3.15, 11, 0.45,
         f"분석 키워드: {keywords}", 14, False, LGRAY)
    today = datetime.date.today().strftime("%Y년 %m월 %d일")
    _txt(s, 1.2, 3.65, 11, 0.4, today, 13, False, LGRAY)
    for i, (lbl, val) in enumerate([
        ("블로그 수집", f"{n_blog:,}건"),
        ("쇼핑 수집",   f"{n_shop:,}개"),
        ("Powered by",  "Naver + Claude"),
    ]):
        x = 1.2 + i * 4.0
        _rect(s, x, 5.5, 3.6, 1.2, RGBColor(0x12,0x25,0x42))
        _rect(s, x, 5.5, 3.6, 0.05, CYAN)
        _txt(s, x+0.15, 5.62, 3.3, 0.35, lbl, 11, False, LGRAY)
        _txt(s, x+0.15, 5.98, 3.3, 0.5, val, 20, True, WHITE)

# ── 슬라이드 2: 블로그 점수
def slide_blog_score(prs, score_data, type_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s); _header(s, "블로그 트렌드 — 관련성 점수 분석",
                    "음료 관련 블로그 필터링 후 키워드별 품질 평가")

    # 좌: 점수 바
    _txt(s, 0.3, 1.15, 5.8, 0.35, "키워드별 평균 관련성 점수", 12, True, CYAN)
    max_s = max((v for _, v, _ in score_data), default=1)
    for i, (kw, score, cnt) in enumerate(score_data[:7]):
        y = 1.55 + i * 0.68
        bw = (score / max(max_s,1)) * 5.5
        _rect(s, 0.3, y+0.28, 5.5, 0.28, RGBColor(0x12,0x25,0x42))
        _rect(s, 0.3, y+0.28, max(bw,0.06), 0.28, CYAN2)
        _txt(s, 0.32, y+0.02, 4.0, 0.26, kw[:18], 11, False, WHITE)
        _txt(s, 4.5,  y+0.02, 1.3, 0.26,
             f"{score:.0f}pt ({cnt}건)", 10, True, CYAN, PP_ALIGN.RIGHT)

    # 우: 콘텐츠 유형 가이드
    _txt(s, 7.0, 1.15, 5.8, 0.35, "콘텐츠 유형별 비중", 12, True, CYAN)
    types_guide = [
        ("🆕 신제품/출시", "+30점"),("☕ 카페메뉴","+25점"),
        ("📈 소비트렌드","+20점"),("💬 구매후기","+20점"),
        ("🔬 성분분석","+15점"),("🛒 구매채널","+15점"),
    ]
    for i, (typ, pt) in enumerate(types_guide):
        y = 1.55 + i * 0.77
        _rect(s, 7.0, y, 5.8, 0.65, RGBColor(0x12,0x25,0x42))
        _rect(s, 7.0, y, 0.06, 0.65, CYAN)
        _txt(s, 7.12, y+0.06, 4.0, 0.3, typ, 12, True, WHITE)
        _txt(s, 11.0, y+0.15, 0.7, 0.3, pt, 12, True, CYAN, PP_ALIGN.RIGHT)
    _footer(s)

# ── 슬라이드 3: TOP 블로그
def slide_top_posts(prs, top_posts):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s); _header(s, "블로그 TOP 10 — 고관련성 게시글",
                    "관련성 점수 기준 상위 10개 음료 관련 게시글")
    # 헤더 행
    _rect(s, 0.2, 1.1, 12.9, 0.38, RGBColor(0x0A,0x1A,0x30))
    _rect(s, 0.2, 1.1, 12.9, 0.04, CYAN)
    for j, (h, x, w) in enumerate([
        ("순위",0.2,0.7),("점수",0.95,0.7),("유형",1.7,1.3),
        ("키워드",3.05,1.4),("제목",4.5,8.55)
    ]):
        _txt(s, x+0.05, 1.12, w-0.1, 0.3, h, 10, True, CYAN)
    medals = ["🥇","🥈","🥉","④","⑤","⑥","⑦","⑧","⑨","⑩"]
    for i, r in enumerate(top_posts[:10]):
        y = 1.5 + i * 0.48
        bg = RGBColor(0x10,0x22,0x3A) if i%2==0 else DGRAY
        _rect(s, 0.2, y, 12.9, 0.45, bg)
        vals = [
            (medals[i] if i<10 else str(i+1), 0.2, 0.7, 10, WHITE),
            (f"{int(r.get('관련성점수',0))}", 0.95, 0.7, 10, CYAN),
            (str(r.get('콘텐츠유형',''))[:9], 1.7, 1.3, 9, WHITE),
            (str(r.get('키워드',''))[:10], 3.05, 1.4, 9, WHITE),
            (str(r.get('제목',''))[:44], 4.5, 8.45, 9, WHITE),
        ]
        for txt, x, w, sz, col in vals:
            _txt(s, x+0.05, y+0.08, w-0.1, 0.3, txt, sz, False, col)
    _footer(s)

# ── 슬라이드 4: 쇼핑
def slide_shop(prs, type_data, price_data):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s); _header(s, "쇼핑 트렌드 — 상품 유형 & 가격 분석",
                    "RTD음료 중심 가격 구조 및 시장 분포 (이상값 제거 후)")

    # 좌: 상품유형 분포
    _txt(s, 0.3, 1.15, 5.8, 0.35, "상품 유형 분포", 12, True, CYAN)
    total = sum(v for _, v in type_data) if type_data else 1
    bar_colors = [CYAN2, RGBColor(0x00,0x96,0xB4),
                  RGBColor(0x00,0x7A,0x96), RGBColor(0x00,0x5F,0x78),
                  RGBColor(0x00,0x48,0x60)]
    for i, (typ, cnt) in enumerate(type_data[:5]):
        y   = 1.58 + i * 0.98
        pct = cnt / total * 100
        bw  = pct / 100 * 5.4
        _txt(s, 0.3, y,      5.4, 0.3, typ, 12, True, WHITE)
        _rect(s, 0.3, y+0.34, 5.4, 0.26, RGBColor(0x12,0x25,0x42))
        _rect(s, 0.3, y+0.34, max(bw,0.06), 0.26,
              bar_colors[i % len(bar_colors)])
        _txt(s, 5.0, y, 0.7, 0.3, f"{pct:.0f}%", 12, True, CYAN, PP_ALIGN.RIGHT)

    # 우: 가격 테이블
    _txt(s, 7.0, 1.15, 6.0, 0.35, "키워드별 RTD 가격 분석", 12, True, CYAN)
    _rect(s, 7.0, 1.55, 6.0, 0.36, RGBColor(0x0A,0x1A,0x30))
    _rect(s, 7.0, 1.55, 6.0, 0.04, CYAN)
    for j, (h, x, w) in enumerate([
        ("키워드",7.0,2.1),("평균 개당",9.15,1.9),("100ml당",11.1,1.85)
    ]):
        _txt(s, x+0.1, 1.57, w-0.15, 0.28, h, 10, True, CYAN)
    for i, (kw, up, mp) in enumerate(price_data[:7]):
        y = 1.94 + i * 0.58
        bg = RGBColor(0x10,0x22,0x3A) if i%2==0 else DGRAY
        _rect(s, 7.0, y, 6.0, 0.54, bg)
        _txt(s, 7.1,  y+0.1, 2.0, 0.34, str(kw)[:12], 11, False, WHITE)
        _txt(s, 9.25, y+0.1, 1.8, 0.34,
             f"{up:,.0f}원" if up else "-", 11, True, CYAN)
        _txt(s, 11.2, y+0.1, 1.7, 0.34,
             f"{mp:,.0f}원" if mp else "-", 11, False, WHITE)
    _footer(s)

# ── 슬라이드 5: NPD 아이디어
def slide_npd(prs, ideas):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s); _header(s, "NPD 인사이트 — 신제품 아이디어 3선",
                    "블로그·쇼핑 데이터 기반 Claude AI 분석 결과")
    card_c = [CYAN2, RGBColor(0x00,0x7A,0x8A), RGBColor(0x00,0x5F,0x6E)]
    for i, idea in enumerate(ideas[:3]):
        x = 0.35 + i * 4.32
        _rect(s, x, 1.2, 4.0, 5.55, DGRAY)
        _rect(s, x, 1.2, 4.0, 0.07, card_c[i%3])
        # 번호
        _rect(s, x+0.15, 1.35, 0.5, 0.5, card_c[i%3])
        _txt(s, x+0.15, 1.35, 0.5, 0.5,
             str(i+1), 18, True, WHITE, PP_ALIGN.CENTER)
        _txt(s, x+0.72, 1.38, 3.15, 0.44,
             str(idea.get('idea',''))[:20], 14, True, WHITE)
        # 구분선 대신 얇은 rect
        _rect(s, x+0.15, 1.92, 3.7, 0.03, LGRAY)
        rows = [
            ("🍹 플레이버", idea.get('flavor','-')),
            ("🏷️ 컨셉",    idea.get('concept','-')),
            ("💰 가격대",   idea.get('price_range','-')),
        ]
        for j, (lbl, val) in enumerate(rows):
            y = 2.05 + j * 0.78
            _txt(s, x+0.15, y,      3.7, 0.28, lbl, 10, False, LGRAY)
            _txt(s, x+0.15, y+0.28, 3.7, 0.38,
                 str(val)[:20], 13, True, CYAN)
        _rect(s, x+0.15, 4.47, 3.7, 0.03, LGRAY)
        _txt(s, x+0.15, 4.55, 3.7, 0.28, "📌 근거", 10, False, LGRAY)
        _txt(s, x+0.15, 4.83, 3.7, 1.7,
             str(idea.get('reason',''))[:90], 10, False, WHITE)
    _footer(s)

# ── 슬라이드 6: 클로징
def slide_closing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _rect(s, 0, 0, 0.45, 7.5, DGRAY)
    _rect(s, 0.45, 3.05, 12.88, 0.05, CYAN)
    _txt(s, 1.5, 1.8, 10, 0.6, "AI NPD SUITE", 18, True, CYAN)
    _txt(s, 1.5, 2.45, 10, 1.0,
         "데이터 기반 음료 트렌드\n분석 완료", 30, True, WHITE)
    _txt(s, 1.5, 3.85, 10, 0.45,
         "Powered by Naver API · Google Gemini · Anthropic Claude",
         13, False, LGRAY)
    _txt(s, 1.5, 4.35, 10, 0.4,
         datetime.date.today().strftime("%Y년 %m월 %d일"), 13, False, LGRAY)
    _footer(s)

# ── 메인 빌드
def build_ppt(blog_df, shop_df, claude_result=None):
    import pandas as pd
    prs = _prs()
    kws = ", ".join(blog_df["키워드"].unique()[:5]) if not blog_df.empty else "-"
    slide_title(prs, kws, len(blog_df), len(shop_df))

    if not blog_df.empty:
        score_data = (blog_df.groupby("키워드")["관련성점수"]
                      .agg(score="mean", cnt="count")
                      .reset_index()
                      .sort_values("score", ascending=False))
        sl = [(r["키워드"], r["score"], int(r["cnt"]))
              for _, r in score_data.iterrows()]
        type_data = list(blog_df["콘텐츠유형"].value_counts().items())
        slide_blog_score(prs, sl, type_data)
        top = blog_df.nlargest(10,"관련성점수")[
            ["관련성점수","콘텐츠유형","키워드","제목"]].to_dict("records")
        slide_top_posts(prs, top)

    if not shop_df.empty:
        td = list(shop_df["상품유형"].value_counts().items())
        rtd = shop_df[shop_df["상품유형"]=="RTD음료"]
        pd_list = []
        for kw in shop_df["키워드"].unique():
            sub = rtd[rtd["키워드"]==kw]
            u = sub["개당가격(원)"].mean() if not sub.empty else None
            mv = sub["100ml당가격(원)"].dropna()
            m = mv.mean() if not mv.empty else None
            pd_list.append((kw, u, m))
        slide_shop(prs, td, pd_list)

    if claude_result and "npd_ideas" in claude_result:
        slide_npd(prs, claude_result["npd_ideas"])

    slide_closing(prs)
    buf = io.BytesIO()
    prs.save(buf); buf.seek(0)
    return buf.read()


# ══════════════════════════════════════════════════════
# ── 앱 소개 + 개발 과정 PPT 생성 ─────────────────────
# ══════════════════════════════════════════════════════

PHASE_DATA = [
    ("PHASE 1", "Instagram 크롤러 시도 → 실패",
     "e53935",
     "• Selenium+BS4로 해시태그 크롤링 시도\n"
     "• ChromeDriver 경로 오류, webdriver_manager 미설치\n"
     "• instaloader 교체 → Instagram 403 차단 (login_required)\n"
     "• 로컬 로그인도 Rate limit → Instagram 크롤링 포기"),

    ("PHASE 2", "네이버 API 전환 · Streamlit Cloud 배포",
     "1e88e5",
     "• 네이버 블로그 + 쇼핑 검색 API 채택\n"
     "• Streamlit Cloud 배포 (requirements.txt, Secrets 관리)\n"
     "• packages.txt 충돌(libgconf), pywin32 오류 해결\n"
     "• API 키 st.secrets로 자동 주입"),

    ("PHASE 3", "데이터 품질 개선",
     "43a047",
     "• 노이즈 필터링 (여행·의료·커피·전자담배 자동 제외)\n"
     "• 관련성 점수 체계 (신제품+30 / 카페메뉴+25 / 트렌드+20)\n"
     "• 쇼핑 상품유형 자동분류 (RTD/농축/건강기능/대용량)\n"
     "• 100ml당 가격 IQR 이상값 자동 제거"),

    ("PHASE 4", "시각화 고도화 (Plotly)",
     "fb8c00",
     "• st.bar_chart → Plotly 전환 (인터랙티브 호버·줌)\n"
     "• 블로그 관련성점수 가로바 + 콘텐츠유형 파이차트\n"
     "• 쇼핑 브랜드 포지셔닝 맵 (가격×인기순위×버블)\n"
     "• 키워드×콘텐츠유형 영향력 지수 히트맵"),

    ("PHASE 5", "AI 기능 추가 (Gemini + Claude)",
     "8e24aa",
     "• Gemini 2.5 Flash → 키워드 10개 추천 (조건별 멀티셀렉트)\n"
     "• 2단계 AI 조건 생성: 카테고리→트렌드·타겟→키워드\n"
     "• Claude Sonnet → 심층 분석 리포트 (JSON 구조화)\n"
     "• NPD 아이디어 3선 카드 렌더링"),

    ("PHASE 6", "DataLab 실검색 지표 연동",
     "00897b",
     "• 네이버 DataLab 검색어 트렌드 API 연동\n"
     "• 월별(3·6개월) / 주별(3·6개월) 탭 전환\n"
     "• 카테고리별 대표 키워드 자동 조회 → Gemini 프롬프트 주입\n"
     "• 네이버 뉴스 헤드라인도 함께 분석에 반영"),

    ("PHASE 7", "PPT + HTML 리포트 출력",
     "00acc1",
     "• python-pptx로 BCG/Deloitte 스타일 PPT 자동 생성 (6슬라이드)\n"
     "• HTML 리포트 (네이비 디자인, 링크 클릭 가능)\n"
     "• session_state 캐시로 버튼 클릭 후에도 다운로드 유지\n"
     "• PPT + HTML 동시 생성·다운로드 지원"),

    ("PHASE 8", "버그픽스 · 안정화",
     "f4511e",
     "• DataLab 라디오 버튼 클릭 시 화면 꺼짐 → if btn 밖 이동\n"
     "• fetch_datalab_range 내부함수 → 최상위함수 이동\n"
     "• 구글 트렌드(pytrends) 403 차단 확인 → 제거\n"
     "• Gemini 2.0-flash 지원종료 → 2.5-flash 교체\n"
     "• CSV·다운로드 버튼 session_state 분리 (화면 리셋 방지)\n"
     "• CSV 한글 인코딩 utf-8-sig 적용"),
]

FEATURE_DATA = [
    ("📝 블로그 트렌드",
     "1b5e20",
     "• 노이즈 자동 필터링\n"
     "• 관련성 점수 체계\n"
     "• 콘텐츠유형 영향력 히트맵\n"
     "• 플레이버 가중 점수 랭킹\n"
     "• TOP 15 게시글 테이블"),
    ("🛒 쇼핑 트렌드",
     "bf360c",
     "• 상품유형 자동분류\n"
     "• IQR 이상값 제거\n"
     "• 개당·100ml당 가격\n"
     "• 브랜드 포지셔닝 맵\n"
     "• 플레이버별 가격 분석"),
    ("🤖 AI · 리포트",
     "0d47a1",
     "• DataLab+뉴스→Gemini 트렌드 추천\n"
     "• DataLab 월별·주별 차트\n"
     "• Claude 심층 분석 카드\n"
     "• NPD 아이디어 3선\n"
     "• PPT + HTML 다운로드"),
]

def _hex(h):
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def slide_intro_cover(prs):
    """슬라이드 1 — 앱 소개 타이틀"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _rect(s, 0, 0, 0.45, 7.5, DGRAY)
    _rect(s, 0.45, 3.2, 12.88, 0.06, CYAN)
    _txt(s, 1.2, 1.0, 11, 0.7, "AI NPD SUITE", 18, True, CYAN)
    _txt(s, 1.2, 1.7, 11, 1.2,
         "앱 소개 및\n바이브 코딩 개발 과정", 36, True, WHITE)
    _txt(s, 1.2, 3.1, 11, 0.5,
         "Claude와의 대화형 개발 기록 — Streamlit Cloud 배포",
         13, False, LGRAY)
    today = datetime.date.today().strftime("%Y년 %m월 %d일")
    _txt(s, 1.2, 3.6, 11, 0.4, today, 12, False, LGRAY)

    specs = [
        ("검색 API", "Naver 블로그·쇼핑·DataLab·뉴스"),
        ("AI 모델", "Gemini 2.5 Flash + Claude Sonnet 4"),
        ("배포", "Streamlit Cloud"),
        ("출력", "PPT (BCG스타일) + HTML 리포트"),
    ]
    for i, (lbl, val) in enumerate(specs):
        x = 1.2 + i * 3.0
        _rect(s, x, 4.5, 2.7, 1.1, DGRAY)
        _rect(s, x, 4.5, 2.7, 0.05, CYAN)
        _txt(s, x+0.1, 4.58, 2.5, 0.3, lbl, 10, False, LGRAY)
        _txt(s, x+0.1, 4.9,  2.5, 0.5, val, 11, True, WHITE)
    _footer(s)


def slide_phase_overview(prs):
    """슬라이드 2 — 개발 단계 개요 (전체 8단계 타임라인)"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "바이브 코딩 개발 과정 — 8단계 타임라인",
            "Claude와 대화형으로 개발한 AI NPD SUITE 구축 여정")

    cols = 4
    for i, (phase, title, color_hex, _) in enumerate(PHASE_DATA):
        row = i // cols
        col = i % cols
        x = 0.25 + col * 3.25
        y = 1.2  + row * 2.85
        c = _hex(color_hex)

        _rect(s, x, y, 3.0, 2.55, DGRAY)
        _rect(s, x, y, 3.0, 0.06, c)

        # 단계 번호 뱃지
        _rect(s, x+0.1, y+0.12, 0.55, 0.42, c)
        _txt(s, x+0.1, y+0.12, 0.55, 0.42,
             str(i+1), 14, True, WHITE, PP_ALIGN.CENTER)

        _txt(s, x+0.72, y+0.14, 2.1, 0.38,
             phase, 10, True, RGBColor(
                 int(color_hex[0:2],16),
                 int(color_hex[2:4],16),
                 int(color_hex[4:6],16)))
        _txt(s, x+0.1,  y+0.6,  2.8, 0.55,
             title, 11, True, WHITE)
        _txt(s, x+0.1,  y+1.18, 2.8, 1.25,
             _.split('\n')[0].lstrip('• '), 9, False, LGRAY)
    _footer(s)


def slide_phase_detail(prs, phase_idx):
    """슬라이드 3~6 — 각 단계 상세 (2단계씩 한 슬라이드)"""
    phase, title, color_hex, detail = PHASE_DATA[phase_idx]
    phase2_idx = phase_idx + 1
    has_second = phase2_idx < len(PHASE_DATA)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)

    if has_second:
        p2, t2, c2_hex, d2 = PHASE_DATA[phase2_idx]
        _header(s, f"{phase}  +  {p2} — 상세",
                f"{title}  /  {t2}")
    else:
        _header(s, f"{phase} — 상세", title)

    def _draw_phase(sx, ph, ttl, chex, det):
        c = _hex(chex)
        _rect(s, sx, 1.15, 6.2, 5.85, DGRAY)
        _rect(s, sx, 1.15, 6.2, 0.07, c)
        _rect(s, sx+0.15, 1.3, 0.7, 0.55, c)
        _txt(s, sx+0.15, 1.3, 0.7, 0.55,
             ph[-1], 18, True, WHITE, PP_ALIGN.CENTER)
        _txt(s, sx+0.95, 1.35, 5.1, 0.5,
             ttl, 15, True, WHITE)
        lines = [l.lstrip('• ') for l in det.strip().split('\n') if l.strip()]
        for j, line in enumerate(lines[:7]):
            y = 2.0 + j * 0.68
            _rect(s, sx+0.15, y+0.12, 0.08, 0.35, c)
            _txt(s, sx+0.3, y+0.06, 5.75, 0.52, line, 11, False, WHITE)

    _draw_phase(0.25, phase, title, color_hex, detail)
    if has_second:
        _draw_phase(7.0, p2, t2, c2_hex, d2)
    _footer(s)


def slide_features(prs):
    """슬라이드 — 현재 기능 구성"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _header(s, "현재 기능 구성",
            "AI NPD SUITE — 블로그·쇼핑·DataLab·AI 통합 인텔리전스 플랫폼")

    for i, (title, chex, detail) in enumerate(FEATURE_DATA):
        x = 0.4 + i * 4.3
        c = _hex(chex)
        _rect(s, x, 1.15, 4.0, 5.85, DGRAY)
        _rect(s, x, 1.15, 4.0, 0.07, c)
        _txt(s, x+0.15, 1.28, 3.7, 0.5, title, 14, True, WHITE)
        lines = [l.lstrip('• ') for l in detail.strip().split('\n') if l.strip()]
        for j, line in enumerate(lines):
            y = 1.9 + j * 0.82
            _rect(s, x+0.15, y+0.16, 0.08, 0.38, c)
            _txt(s, x+0.3, y+0.08, 3.55, 0.62, line, 11, False, WHITE)
    _footer(s)


def slide_intro_closing(prs):
    """마지막 슬라이드 — 클로징"""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s)
    _rect(s, 0, 0, 0.45, 7.5, DGRAY)
    _rect(s, 0.45, 3.1, 12.88, 0.05, CYAN)
    _txt(s, 1.2, 1.6, 10, 0.6, "AI NPD SUITE", 18, True, CYAN)
    _txt(s, 1.2, 2.3, 11, 1.0,
         "바이브 코딩으로 만든\n음료 트렌드 분석기", 30, True, WHITE)
    _txt(s, 1.2, 3.8, 11, 0.45,
         "Claude × Natural Lab — 데이터 기반 NPD 인텔리전스",
         13, False, LGRAY)
    _txt(s, 1.2, 4.3, 11, 0.4,
         datetime.date.today().strftime("%Y년 %m월 %d일"), 12, False, LGRAY)
    _footer(s)


def build_intro_ppt():
    """앱 소개 + 개발 과정 PPT 빌드 → bytes 반환"""
    prs = _prs()

    # 슬라이드 1: 커버
    slide_intro_cover(prs)

    # 슬라이드 2: 전체 타임라인
    slide_phase_overview(prs)

    # 슬라이드 3~6: 2단계씩 상세 (4슬라이드)
    for idx in range(0, len(PHASE_DATA), 2):
        slide_phase_detail(prs, idx)

    # 슬라이드 7: 기능 구성
    slide_features(prs)

    # 슬라이드 8: 클로징
    slide_intro_closing(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
